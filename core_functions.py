from UI_design import *
from PyQt5.QtGui import QIcon, QPixmap
from pptx import Presentation
from pptx.util import Pt
from collections import deque
from pptx.enum.text import MSO_AUTO_SIZE
from googleSheetOperator import *
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import datetime
import math
import os

class core_functions(Ui_MainWindow):
    songCountList = ["3", "4", "5", "6", "7", "8"]
    currentSongCountChoice = 0
    snapshotFolder = "../PowerPoint_Templates/snapshots"
    templateFolder = "../PowerPoint_Templates"
    singleSongsFolder = "../PowerPoint_Single_Songs"
    outputFolder = "../Output"
    templateList = []
    singleSongList = []
    songSelectionText = None

    def __init__(self):
        print("core functions initialization")
        super().__init__()

    def setupOtherAttributes(self):

        print("Class name for setupOtherAttributes: " + self.__class__.__name__)

        self.centralwidget.setFixedSize(696, 443)
        self.dateEdit.setDisplayFormat("yyyy/MM/dd")
        self.dateEdit.setDate(QtCore.QDate.currentDate())
        self.setUpEvents()

    def setUpEvents(self):
        self.proceedButton.clicked.connect(self.process)
        self.songCountComboBox.addItems(self.songCountList)
        self.songCountComboBox.currentIndexChanged.connect(self.songCountComboBoxCallBack)
        self.templateComboBox.currentIndexChanged.connect(self.templateComboBoxCallBack)

        self.currentSongCountChoice = self.songCountComboBox.currentText()

        for i in range(4, 9):
            exec("self.song_" + str(i) + "_label.setEnabled(False)")
            exec("self.song_" + str(i) + "_comboBox.setEnabled(False)")

        # self.preview_image.setPixmap(QPixmap(self.snapshotFolder+"/ebcsv_2012_nov_dec_template_ppt.png"))

        # os.listdir(self.snapshotFolder)[0]
        # print(self.findFirstImage(self.snapshotFolder))
        self.preview_image.setPixmap(QPixmap(self.findFirstImage(self.snapshotFolder)))
        self.preview_image.setScaledContents(True)

        if not os.path.exists(self.outputFolder):
            os.makedirs(self.outputFolder)

        self.setUpTemplateList()
        self.setUpSongList()
        self.setUpSongDropDowns()

        print("completed")

    def setUpSongDropDowns(self):

        for i in range(1, 9):
            exec("self.song_" + str(i) + "_comboBox.addItems(self.singleSongList)")
            exec("self.song_" + str(i) + "_comboBox.setCurrentIndex(-1)")

    def setUpSongList(self):

        dirs = os.listdir(self.singleSongsFolder)
        for file in dirs:
            if (".pptx" in file):
                self.singleSongList.append(file)

    def findFirstImage(self, folderPath):
        fileName = ""

        dirs = os.listdir(folderPath)

        for f in dirs:
            if ".png" in f:
                fileName = f
                break

        return folderPath + "/" + fileName

    def templateComboBoxCallBack(self):

        imageText = self.templateComboBox.currentText()
        self.preview_image.setPixmap(QPixmap(self.snapshotFolder + "/" + imageText))
        self.preview_image.setScaledContents(True)

    def setUpTemplateList(self):

        dirs = os.listdir(self.snapshotFolder)
        for file in dirs:
            # if ".ppt" in file or ".pptx" in file:
            if ".png" in file:
                self.templateList.append(file)

        self.templateComboBox.addItems(self.templateList)

    def songCountComboBoxCallBack(self):
        print("comboBox Value has changed!")
        currentChoice = int(self.currentSongCountChoice)
        newChoice = int(self.songCountComboBox.currentText())
        self.currentSongCountChoice = newChoice

        if (currentChoice > newChoice):
            for i in range(currentChoice, newChoice, -1):
                exec("self.song_" + str(i) + "_label.setEnabled(False)")
                exec("self.song_" + str(i) + "_comboBox.setEnabled(False)")
                exec("self.song_" + str(i) + "_comboBox.setCurrentIndex(-1)")

        elif (currentChoice < newChoice):
            for i in range(newChoice, currentChoice, -1):
                exec("self.song_" + str(i) + "_label.setEnabled(True)")
                exec("self.song_" + str(i) + "_comboBox.setEnabled(True)")

        print(str(currentChoice) + " " + str(newChoice))

    def extractText(self, original_slide):
        texts = deque()
        footnoteObject = None
        footnotes = deque()
        i = 0
        for shape in original_slide.shapes:
            print("Shape #" + str(i))
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print("P: " + paragraph.text)
                # if "'" in paragraph.text:
                if (paragraph.text.strip(' \t\n\r') == ''):
                    continue
                temp = "".join(paragraph.text)
                if (i == 2):
                    footnotes.append(temp)
                    #footnotes += temp
                else:
                    texts.append(temp)

                    # for run in paragraph.runs:
                    #     print(run.text)
                    #     if (i == 2 ):
                    #         footnotes += run.text
                    #     else:
                    #         texts.append(run.text)
            i = i + 1

        title = texts.popleft()
        return (title, texts, footnotes)

    def process(self):

        currentDate = self.dateEdit.text()
        dayOfWeek = datetime.datetime.strptime(currentDate, "%Y/%m/%d").strftime('%A')

        songCount = int(self.songCountComboBox.currentText())
        for i in range(1, songCount + 1):
            result = eval("self.song_" + str(i) + "_comboBox.currentText()")
            if (result == ''):
                print("Please select a song")
                return

        T = self.templateFolder + "/" + self.templateComboBox.currentText().replace("_ppt.png", ".ppt").replace(
            "_pptx.png", ".pptx")
        print("master PPTX path: " + T)
        master = Presentation(T)
        print("# slides: " + str(len(master.slides)))

        ### retrieve master slide text format ###

        # slide #0 is the first slide that contains main information
        # actual songs begin at slide #1 onwards.

        # title format
        titleObject = master.slides[1].shapes[0].text_frame.paragraphs[0]  # first run is the title
        print("Title: " + titleObject.text)
        # print(titleObject.font.size)

        # texts format
        textObject = master.slides[1].shapes[1].text_frame.paragraphs[0]  # second run is the content
        print("Text: " + textObject.text)
        # print(textObject.font.size)

        # footNotes (CopyRights info)
        footNotesObject = master.slides[1].shapes[2].text_frame.paragraphs[0]
        print("FootNotes: " + footNotesObject.text)
        # print(footNotesObject.font.size)

        # removing existing slides from master (except the first slide), if any
        for i in range(-1, len(master.slides) - 2):
            print(i)
            rId = master.slides._sldIdLst[i].rId
            master.part.drop_rel(rId)
            del master.slides._sldIdLst[i]

            # -------------------------------------------------------------------------------
        for j in range(1, songCount + 1):
            S = eval("self.singleSongsFolder + \"/\" + self.song_" + str(j) + "_comboBox.currentText()")
            print("source path: " + S)
            source = Presentation(S)

            # print(len(source1.slides))
            for slide in source.slides:

                (title, texts, footNotes) = self.extractText(slide)
                print("title: " + title)
                print("texts: " + texts.__str__())
                print("footNotes: " + footNotes.__str__())

                no_paragraphs = len(texts)

                subTexts = deque()

                if (dayOfWeek == "Thursday"):
                    no_sub_slides = math.ceil(no_paragraphs / 2)
                    max_paragraphs_per_slide = 2

                    for L in range(0, no_sub_slides):
                        T = deque()
                        for y in range(0, max_paragraphs_per_slide):
                            if texts:
                                T.append(texts.popleft())

                        subTexts.append(T)

                else:
                    no_sub_slides = 1
                    max_paragraphs_per_slide = no_paragraphs
                    subTexts = texts

                for ss in range(0, no_sub_slides):

                    # assuming this is the master layout at index= 0 for the given pptx file.
                    s = master.slides.add_slide(master.slide_layouts[0])

                    self.addTitle(s, titleObject, title)
                    # tf = s.shapes[0].text_frame
                    # tf.clear()
                    # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    # p = tf.paragraphs[0]
                    # run = p.add_run()
                    # self.copyFontStyle(run, titleObject)
                    # run.text = title

                    self.addBody(s, textObject, no_sub_slides, subTexts, ss)
                    # body_shape = s.placeholders[1]
                    # btf = body_shape.text_frame
                    # btf.clear()
                    # btf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    # bp = btf.paragraphs[0]
                    # bp.level = 1
                    # brun = bp.add_run()
                    # self.copyFontStyle(brun, textObject)
                    # if (no_sub_slides == 1):
                    #     for lyrics in subTexts:
                    #         brun.text += lyrics + "\n"
                    # else:
                    #     for lyrics in subTexts[ss]:
                    #         brun.text += lyrics + "\n"

                    if (ss == no_sub_slides - 1):
                        if (len(footNotes) > 0):
                            self.addFootNotes(s, footNotesObject, footNotes)
                            # print("found footNotes: " + footNotes)
                            # fntf = s.shapes[2].text_frame
                            # fntf.clear()
                            # fntf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                            # fnp = fntf.paragraphs[0]
                            # fnrun = fnp.add_run()
                            # self.copyFontStyle(fnrun, footNotesObject)
                            # for ly in footNotes:
                            #     fnrun.text += ly
# -------------------------------------------------------------------------------
        # Adding announcements
        (announcements, hasResults) = bulletin().retrieveAnnouncements(datetime.datetime.strptime(currentDate, '%Y/%m/%d').strftime('%m/%d/%Y').replace("/0", "/").lstrip("0"))
        if hasResults:
            # print("announcement #1: " + announcements[0])
            # print("announcement #2: " + announcements[1])
            # print("announcement #3: " + announcements[2])

            no_announcements = 0
            for temp in announcements:
                if temp != '':
                    no_announcements += 1

            for k in range(0, no_announcements):
                if announcements[k] != '':
                    s_ = master.slides.add_slide(master.slide_layouts[0])
                    self.addTitle(s_, titleObject, "Announcements (" + str(k+1) + " of " + str(no_announcements) +")")
                    self.addBody(s_, textObject, 0, announcements[k], 0, True)
        else:
            print("selected Date is not found!")

        master.save(self.outputFolder + "/" + "ebcsv_" + dayOfWeek + "_" + currentDate.replace("/", "_").replace("_0", "_") + ".pptx")

    def addTitle(self, destination, titleObject, title):
        print("=== Adding title ===")
        tf = destination.shapes[0].text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tf.paragraphs[0]
        run = p.add_run()
        self.copyFontStyle(run, titleObject)
        run.text = title

    def addBody(self, destination, textObject, no_sub_slides, subTexts, ss, isAnnouncement = False):
        print("=== Adding body ===")
        body_shape = destination.placeholders[1]
        btf = body_shape.text_frame
        btf.clear()
        btf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        bp = btf.paragraphs[0]
        #bp.level = 1
        brun = bp.add_run()
        self.copyFontStyle(brun, textObject)
        if not isAnnouncement:
            if (no_sub_slides == 1):
                for lyrics in subTexts:
                    brun.text += lyrics + "\n"
            else:
                for lyrics in subTexts[ss]:
                    brun.text += lyrics + "\n"
        else:
            bp.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            brun.text = subTexts

    def addFootNotes(self, destination, footNotesObject, footNotes):
        print("=== Adding footNotes ===")
        print("found footNotes: " + footNotes.__str__())
        fntf = destination.shapes[2].text_frame
        fntf.clear()
        fntf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        fnp = fntf.paragraphs[0]
        fnrun = fnp.add_run()
        self.copyFontStyle(fnrun, footNotesObject)
        for ly in footNotes:
                fnrun.text += ly + "\n"

    def getFile(self, fileName, folderPath):

        dirs = os.listdir(folderPath)
        for file in dirs:
            # print("file: " + file + ", fileName: " + fileName)
            if file == fileName:
                return file

    def copyFontStyle(self, destinationObj, sourceObj):
        destinationObj.font.size = sourceObj.font.size  # setting font size  based on master layout
        destinationObj.font.bold = sourceObj.font.bold  # setting font style based on master layout
        destinationObj.font.name = sourceObj.font.name  # setting name size  based on master layout
        destinationObj.font.italic = sourceObj.font.italic
